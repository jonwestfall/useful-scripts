#!/usr/bin/env python3
"""
anonymize_student_work.py

Create anonymized copies of student submissions and a private crosswalk mapping
anonymous IDs back to original files / identifiers.

Supported:
  - PDF: .pdf, using PyMuPDF true redaction annotations
  - Word: .docx, using python-docx
  - PowerPoint: .pptx, using python-pptx

Not supported directly:
  - Legacy .doc / .ppt files. Convert these to .docx/.pptx first, or add a
    LibreOffice conversion step before running this script.

Install:
  python3 -m venv .venv
  source .venv/bin/activate
  pip install pymupdf python-docx python-pptx

Basic usage:
  python anonymize_student_work.py ./submissions --output ./anonymized

With roster/name list:
  python anonymize_student_work.py ./submissions --output ./anonymized --names roster.csv

Dry run:
  python anonymize_student_work.py ./submissions --output ./anonymized --names roster.csv --dry-run

Roster formats:
  - Plain text: one student name or email per line
  - CSV: columns such as name, first, last, email, student_id are recognized.
    If no recognized columns are found, every non-empty cell is treated as a
    possible identifier.

Important limitations:
  - Scanned PDFs/images cannot be text-redacted without OCR. This script flags
    low/no-text PDFs in the manifest notes.
  - Names embedded in screenshots/images inside Word, PowerPoint, or PDF files
    are not removed.
  - Comments, tracked changes, embedded objects, and unusual custom metadata may
    still contain identifying information. Use the generated manifest notes and
    spot-check outputs before sharing for blind review.
"""

from __future__ import annotations

import argparse
import csv
import datetime as dt
import hashlib
import os
import re
import shutil
import sys
import zipfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Iterable, Optional

SUPPORTED_EXTS = {".pdf", ".docx", ".pptx"}
LEGACY_OFFICE_EXTS = {".doc", ".ppt"}


# -----------------------------
# Data structures
# -----------------------------

@dataclass
class StudentRecord:
    raw_values: set[str] = field(default_factory=set)

    def add(self, value: Optional[str]) -> None:
        if value is None:
            return
        value = str(value).strip()
        if value:
            self.raw_values.add(value)

    @property
    def terms(self) -> set[str]:
        """Generate useful search/redaction terms from roster values."""
        terms: set[str] = set()

        for value in self.raw_values:
            v = value.strip()
            if not v:
                continue

            terms.add(v)

            # Email variants
            if "@" in v:
                local = v.split("@", 1)[0]
                terms.add(local)
                terms.add(local.replace(".", " "))
                terms.add(local.replace("_", " "))
                terms.add(local.replace("-", " "))

            # Name variants
            # "Last, First" -> "First Last"
            if "," in v:
                pieces = [p.strip() for p in v.split(",", 1)]
                if len(pieces) == 2 and pieces[0] and pieces[1]:
                    terms.add(f"{pieces[1]} {pieces[0]}")
                    terms.add(f"{pieces[0]}, {pieces[1]}")

            tokens = re.findall(r"[A-Za-z][A-Za-z'\-]+", v)
            if len(tokens) >= 2:
                first = tokens[0]
                last = tokens[-1]
                terms.add(f"{first} {last}")
                terms.add(f"{last}, {first}")
                terms.add(f"{first[0]}. {last}")
                terms.add(f"{first[0]} {last}")

        # Avoid dangerous, overbroad one-character terms.
        return {t.strip() for t in terms if len(t.strip()) >= 2}


@dataclass
class FileResult:
    anon_id: str
    source_file: Path
    output_file: Optional[Path]
    status: str
    identifiers_found: list[str] = field(default_factory=list)
    identifiers_removed: list[str] = field(default_factory=list)
    inferred_identifiers: list[str] = field(default_factory=list)
    metadata_identifiers: list[str] = field(default_factory=list)
    notes: list[str] = field(default_factory=list)


# -----------------------------
# General helpers
# -----------------------------

def natural_sort_key(path: Path) -> list[object]:
    return [int(part) if part.isdigit() else part.lower()
            for part in re.split(r"(\d+)", path.name)]


def safe_filename_id(n: int, prefix: str, width: int) -> str:
    return f"{prefix}{n:0{width}d}"


def sha256_file(path: Path) -> str:
    h = hashlib.sha256()
    with path.open("rb") as f:
        for block in iter(lambda: f.read(1024 * 1024), b""):
            h.update(block)
    return h.hexdigest()


def casefold_unique(values: Iterable[str]) -> list[str]:
    seen = set()
    out = []
    for v in values:
        v = str(v).strip()
        if not v:
            continue
        key = v.casefold()
        if key not in seen:
            seen.add(key)
            out.append(v)
    return out


def redact_text_literal(text: str, terms: Iterable[str], replacement: str) -> tuple[str, list[str]]:
    """Case-insensitive literal replacement. Longer terms first."""
    changed_terms = []
    new_text = text
    for term in sorted(casefold_unique(terms), key=len, reverse=True):
        # Word-ish boundaries for names/IDs, but allow punctuation in emails.
        pattern = re.compile(re.escape(term), re.IGNORECASE)
        if pattern.search(new_text):
            new_text = pattern.sub(replacement, new_text)
            changed_terms.append(term)
    return new_text, changed_terms


def find_terms_in_text(text: str, terms: Iterable[str]) -> list[str]:
    found = []
    for term in sorted(casefold_unique(terms), key=len, reverse=True):
        if re.search(re.escape(term), text, flags=re.IGNORECASE):
            found.append(term)
    return casefold_unique(found)


def infer_identifiers_from_filename(path: Path) -> list[str]:
    """
    Guess identifiers from common student submission filenames.
    Examples:
      Smith_Jane_final.docx
      Jane Smith Research Paper.pdf
      psych101 - Smith, Jane.pptx
    """
    stem = path.stem
    cleaned = re.sub(r"[_\-]+", " ", stem)
    cleaned = re.sub(r"\b(final|draft|paper|essay|assignment|submission|research|project|presentation|ppt|docx|pdf)\b",
                     " ", cleaned, flags=re.IGNORECASE)
    cleaned = re.sub(r"\s+", " ", cleaned).strip()

    candidates = []
    if "," in stem:
        # Look for Last, First anywhere in filename.
        m = re.search(r"([A-Z][A-Za-z'\-]+)\s*,\s*([A-Z][A-Za-z'\-]+)", stem)
        if m:
            candidates.append(f"{m.group(2)} {m.group(1)}")
            candidates.append(f"{m.group(1)}, {m.group(2)}")

    # Look for two adjacent title-case name-like tokens.
    tokens = re.findall(r"\b[A-Z][A-Za-z'\-]{1,}\b", cleaned)
    if len(tokens) >= 2:
        # Use first two and last two, because LMS filenames vary.
        candidates.append(f"{tokens[0]} {tokens[1]}")
        candidates.append(f"{tokens[-2]} {tokens[-1]}")

    return casefold_unique(candidates)


def infer_identifiers_from_text_head(text: str) -> list[str]:
    """
    Guess identifiers from first-page / first-slide / first-paragraph labels.
    Conservative by design.
    """
    head = "\n".join(text.splitlines()[:40])
    candidates = []

    label_patterns = [
        r"\bName\s*:\s*([A-Z][A-Za-z'\-]+(?:\s+[A-Z][A-Za-z'\-]+){0,3})",
        r"\bStudent\s*:\s*([A-Z][A-Za-z'\-]+(?:\s+[A-Z][A-Za-z'\-]+){0,3})",
        r"\bAuthor\s*:\s*([A-Z][A-Za-z'\-]+(?:\s+[A-Z][A-Za-z'\-]+){0,3})",
        r"\bSubmitted\s+by\s*:?\s*([A-Z][A-Za-z'\-]+(?:\s+[A-Z][A-Za-z'\-]+){0,3})",
        r"\bBy\s+([A-Z][A-Za-z'\-]+(?:\s+[A-Z][A-Za-z'\-]+){1,3})",
    ]

    for pattern in label_patterns:
        for m in re.finditer(pattern, head, flags=re.IGNORECASE):
            value = m.group(1).strip()
            # Avoid grabbing course titles like "Introduction To Psychology"
            if len(value.split()) <= 4:
                candidates.append(value)

    email_matches = re.findall(r"\b[A-Z0-9._%+\-]+@[A-Z0-9.\-]+\.[A-Z]{2,}\b", head, flags=re.IGNORECASE)
    candidates.extend(email_matches)

    return casefold_unique(candidates)


def load_roster(path: Optional[Path]) -> dict[str, StudentRecord]:
    """
    Returns a dictionary keyed by a stable internal roster key.
    This script does not require matching a file to one roster row; all roster
    terms are used for redaction, and any terms found in a file are listed.
    """
    records: dict[str, StudentRecord] = {}

    if not path:
        return records

    if not path.exists():
        raise FileNotFoundError(f"Names/roster file not found: {path}")

    def add_record(key: str, values: Iterable[str]) -> None:
        rec = records.setdefault(key, StudentRecord())
        for value in values:
            rec.add(value)

    if path.suffix.lower() == ".csv":
        with path.open("r", encoding="utf-8-sig", newline="") as f:
            sample = f.read(4096)
            f.seek(0)
            has_header = csv.Sniffer().has_header(sample) if sample.strip() else True
            reader = csv.DictReader(f) if has_header else None

            if reader and reader.fieldnames:
                fieldnames = [name.strip() for name in reader.fieldnames]
                normalized = {name: re.sub(r"[^a-z0-9]+", "_", name.lower()).strip("_")
                              for name in fieldnames}
                for idx, row in enumerate(reader, start=1):
                    row_values = {k: (v or "").strip() for k, v in row.items()}
                    rec_values = []

                    # Common explicit columns
                    for original, norm in normalized.items():
                        if norm in {"name", "student_name", "full_name", "preferred_name",
                                    "email", "student_email", "e_mail", "student_id", "id"}:
                            rec_values.append(row_values.get(original, ""))

                    # First + last
                    first = ""
                    last = ""
                    for original, norm in normalized.items():
                        if norm in {"first", "first_name", "given", "given_name"}:
                            first = row_values.get(original, "")
                        if norm in {"last", "last_name", "surname", "family", "family_name"}:
                            last = row_values.get(original, "")
                    if first and last:
                        rec_values.append(f"{first} {last}")
                        rec_values.append(f"{last}, {first}")

                    # Fallback: every non-empty cell if explicit columns failed.
                    if not any(v.strip() for v in rec_values):
                        rec_values = [v for v in row_values.values() if v.strip()]

                    add_record(f"row_{idx}", rec_values)
            else:
                f.seek(0)
                reader2 = csv.reader(f)
                for idx, row in enumerate(reader2, start=1):
                    add_record(f"row_{idx}", [cell.strip() for cell in row if cell.strip()])
    else:
        with path.open("r", encoding="utf-8-sig") as f:
            for idx, line in enumerate(f, start=1):
                line = line.strip()
                if line and not line.startswith("#"):
                    add_record(f"line_{idx}", [line])

    return records


def all_roster_terms(records: dict[str, StudentRecord]) -> list[str]:
    terms: list[str] = []
    for rec in records.values():
        terms.extend(rec.terms)
    return casefold_unique(terms)


def build_terms_for_file(source: Path, extracted_text: str, roster_terms: list[str],
                         metadata_terms: list[str], auto_infer: bool = True) -> tuple[list[str], list[str]]:
    inferred = []
    if auto_infer:
        inferred.extend(infer_identifiers_from_filename(source))
        inferred.extend(infer_identifiers_from_text_head(extracted_text))

    all_terms = []
    all_terms.extend(roster_terms)
    all_terms.extend(metadata_terms)
    all_terms.extend(inferred)

    # Avoid redacting common assignment words accidentally.
    stop_terms = {
        "final paper", "research paper", "psychology", "introduction",
        "assignment", "student", "professor", "course", "submitted"
    }
    all_terms = [t for t in casefold_unique(all_terms)
                 if t.casefold() not in stop_terms and len(t.strip()) >= 2]
    return all_terms, casefold_unique(inferred)


# -----------------------------
# DOCX support
# -----------------------------

def require_docx():
    try:
        import docx  # noqa: F401
    except ImportError as e:
        raise RuntimeError("Missing dependency: python-docx. Install with: pip install python-docx") from e


def iter_docx_paragraphs(document):
    """Yield paragraphs in body, tables, headers, and footers."""
    for p in document.paragraphs:
        yield p

    for table in document.tables:
        for row in table.rows:
            for cell in row.cells:
                for p in cell.paragraphs:
                    yield p
                for nested in cell.tables:
                    for row2 in nested.rows:
                        for cell2 in row2.cells:
                            for p in cell2.paragraphs:
                                yield p

    for section in document.sections:
        for part in [section.header, section.footer,
                     section.first_page_header, section.first_page_footer,
                     section.even_page_header, section.even_page_footer]:
            for p in part.paragraphs:
                yield p
            for table in part.tables:
                for row in table.rows:
                    for cell in row.cells:
                        for p in cell.paragraphs:
                            yield p


def replace_paragraph_runs(paragraph, terms: list[str], replacement: str) -> list[str]:
    """
    Replace across runs by collapsing changed paragraph text into the first run.
    This may simplify character-level formatting in the changed paragraph, but
    catches identifiers split across runs.
    """
    if not paragraph.runs:
        return []

    original = "".join(run.text for run in paragraph.runs)
    new_text, changed = redact_text_literal(original, terms, replacement)
    if changed:
        paragraph.runs[0].text = new_text
        for run in paragraph.runs[1:]:
            run.text = ""
    return changed


def get_docx_text_and_metadata(path: Path) -> tuple[str, list[str]]:
    require_docx()
    from docx import Document

    document = Document(path)
    texts = []
    for p in iter_docx_paragraphs(document):
        if p.text:
            texts.append(p.text)

    cp = document.core_properties
    metadata = []
    for attr in ["author", "last_modified_by", "title", "subject", "keywords", "comments", "category"]:
        value = getattr(cp, attr, None)
        if value:
            metadata.append(str(value))

    return "\n".join(texts), casefold_unique(metadata)


def clear_docx_metadata(document, anon_id: str) -> None:
    cp = document.core_properties
    cp.author = "Anonymous"
    cp.last_modified_by = "Anonymous"
    cp.title = f"Anonymous submission {anon_id}"
    cp.subject = ""
    cp.keywords = ""
    cp.comments = ""
    cp.category = ""


def anonymize_docx(source: Path, output: Path, anon_id: str, terms: list[str]) -> tuple[list[str], list[str]]:
    require_docx()
    from docx import Document

    document = Document(source)
    removed = []

    # Add visible anonymous ID at top.
    p = document.paragraphs[0].insert_paragraph_before(f"Anonymous ID: {anon_id}") if document.paragraphs else document.add_paragraph(f"Anonymous ID: {anon_id}")
    try:
        p.style = "Title"
    except Exception:
        pass

    for paragraph in iter_docx_paragraphs(document):
        removed.extend(replace_paragraph_runs(paragraph, terms, "[REDACTED]"))

    clear_docx_metadata(document, anon_id)

    output.parent.mkdir(parents=True, exist_ok=True)
    document.save(output)

    notes = []
    if docx_has_comments_or_tracked_changes(output):
        notes.append("Possible comments or tracked changes detected; manually inspect before sharing.")

    return casefold_unique(removed), notes


def docx_has_comments_or_tracked_changes(path: Path) -> bool:
    """Inspect the zipped OOXML package for common review artifacts."""
    try:
        with zipfile.ZipFile(path) as z:
            names = set(z.namelist())
            if any("comments" in name.lower() for name in names):
                return True
            xml_files = [n for n in names if n.endswith(".xml") and n.startswith("word/")]
            review_tags = [b"<w:ins", b"<w:del", b"<w:commentRangeStart", b"<w:commentReference"]
            for xml_name in xml_files:
                data = z.read(xml_name)
                if any(tag in data for tag in review_tags):
                    return True
    except Exception:
        return True
    return False


# -----------------------------
# PPTX support
# -----------------------------

def require_pptx():
    try:
        import pptx  # noqa: F401
    except ImportError as e:
        raise RuntimeError("Missing dependency: python-pptx. Install with: pip install python-pptx") from e


def iter_pptx_text_paragraphs(prs):
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                for paragraph in shape.text_frame.paragraphs:
                    yield paragraph
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            yield paragraph


def get_pptx_text_and_metadata(path: Path) -> tuple[str, list[str]]:
    require_pptx()
    from pptx import Presentation

    prs = Presentation(path)
    texts = []
    for p in iter_pptx_text_paragraphs(prs):
        if p.text:
            texts.append(p.text)

    cp = prs.core_properties
    metadata = []
    for attr in ["author", "last_modified_by", "title", "subject", "keywords", "comments", "category"]:
        value = getattr(cp, attr, None)
        if value:
            metadata.append(str(value))

    return "\n".join(texts), casefold_unique(metadata)


def replace_pptx_paragraph_runs(paragraph, terms: list[str], replacement: str) -> list[str]:
    if not paragraph.runs:
        return []
    original = "".join(run.text for run in paragraph.runs)
    new_text, changed = redact_text_literal(original, terms, replacement)
    if changed:
        paragraph.runs[0].text = new_text
        for run in paragraph.runs[1:]:
            run.text = ""
    return changed


def clear_pptx_metadata(prs, anon_id: str) -> None:
    cp = prs.core_properties
    cp.author = "Anonymous"
    cp.last_modified_by = "Anonymous"
    cp.title = f"Anonymous submission {anon_id}"
    cp.subject = ""
    cp.keywords = ""
    cp.comments = ""
    cp.category = ""


def anonymize_pptx(source: Path, output: Path, anon_id: str, terms: list[str]) -> tuple[list[str], list[str]]:
    require_pptx()
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation(source)
    removed = []

    # Add visible anonymous ID to first slide.
    if prs.slides:
        slide = prs.slides[0]
        box = slide.shapes.add_textbox(Inches(0.25), Inches(0.15), Inches(4), Inches(0.3))
        tf = box.text_frame
        tf.text = f"Anonymous ID: {anon_id}"
        try:
            tf.paragraphs[0].runs[0].font.size = Pt(14)
        except Exception:
            pass

    for paragraph in iter_pptx_text_paragraphs(prs):
        removed.extend(replace_pptx_paragraph_runs(paragraph, terms, "[REDACTED]"))

    clear_pptx_metadata(prs, anon_id)

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)

    notes = []
    if pptx_has_comments(output):
        notes.append("Possible comments detected; manually inspect before sharing.")

    return casefold_unique(removed), notes


def pptx_has_comments(path: Path) -> bool:
    try:
        with zipfile.ZipFile(path) as z:
            names = set(z.namelist())
            return any("comment" in name.lower() for name in names)
    except Exception:
        return True


# -----------------------------
# PDF support
# -----------------------------

def require_fitz():
    try:
        import fitz  # noqa: F401
    except ImportError as e:
        raise RuntimeError("Missing dependency: PyMuPDF. Install with: pip install pymupdf") from e


def get_pdf_text_and_metadata(path: Path) -> tuple[str, list[str]]:
    require_fitz()
    import fitz

    doc = fitz.open(path)
    texts = []
    for page in doc:
        try:
            texts.append(page.get_text("text") or "")
        except Exception:
            pass

    metadata = []
    for value in (doc.metadata or {}).values():
        if value:
            metadata.append(str(value))

    doc.close()
    return "\n".join(texts), casefold_unique(metadata)


def anonymize_pdf(source: Path, output: Path, anon_id: str, terms: list[str]) -> tuple[list[str], list[str]]:
    require_fitz()
    import fitz

    doc = fitz.open(source)
    removed = []
    notes = []

    total_text_len = 0
    for page in doc:
        page_text = page.get_text("text") or ""
        total_text_len += len(page_text)

        # Add a visible ID in the page margin/top area.
        if page.number == 0:
            page.insert_text((36, 24), f"Anonymous ID: {anon_id}", fontsize=11, color=(0, 0, 0))

        # Redact each literal term where PyMuPDF can locate it.
        for term in sorted(casefold_unique(terms), key=len, reverse=True):
            rects = page.search_for(term)
            if rects:
                removed.append(term)
            for rect in rects:
                page.add_redact_annot(rect, text="[REDACTED]", fill=(1, 1, 1), text_color=(0, 0, 0))

        # Apply redactions after all annotations on this page.
        try:
            page.apply_redactions()
        except Exception as e:
            notes.append(f"Could not apply one or more PDF redactions on page {page.number + 1}: {e}")

    if total_text_len < 50:
        notes.append("Very little extractable text found; PDF may be scanned/image-based and needs OCR/manual redaction.")

    # Strip common PDF metadata.
    doc.set_metadata({
        "title": f"Anonymous submission {anon_id}",
        "author": "Anonymous",
        "subject": "",
        "keywords": "",
        "creator": "",
        "producer": "",
        "creationDate": "",
        "modDate": "",
    })

    output.parent.mkdir(parents=True, exist_ok=True)
    # garbage=4 deflates and removes unused objects where possible.
    doc.save(output, garbage=4, deflate=True, clean=True)
    doc.close()

    return casefold_unique(removed), notes


# -----------------------------
# Processing and manifest
# -----------------------------

def extract_text_and_metadata(path: Path) -> tuple[str, list[str], list[str]]:
    ext = path.suffix.lower()
    notes = []
    if ext == ".pdf":
        text, metadata = get_pdf_text_and_metadata(path)
    elif ext == ".docx":
        text, metadata = get_docx_text_and_metadata(path)
    elif ext == ".pptx":
        text, metadata = get_pptx_text_and_metadata(path)
    else:
        text, metadata = "", []
        if ext in LEGACY_OFFICE_EXTS:
            notes.append("Legacy Office format not processed. Convert to .docx/.pptx first.")
        else:
            notes.append("Unsupported file type.")
    return text, metadata, notes


def process_file(path: Path, anon_id: str, output_dir: Path, roster_terms: list[str],
                 auto_infer: bool, dry_run: bool) -> FileResult:
    ext = path.suffix.lower()
    output_file = output_dir / f"{anon_id}{ext}"

    result = FileResult(
        anon_id=anon_id,
        source_file=path,
        output_file=None if dry_run else output_file,
        status="pending",
    )

    if ext not in SUPPORTED_EXTS:
        result.status = "skipped"
        if ext in LEGACY_OFFICE_EXTS:
            result.notes.append("Legacy Office format skipped. Convert to .docx/.pptx first.")
        else:
            result.notes.append("Unsupported file type.")
        return result

    try:
        text, metadata_terms, extract_notes = extract_text_and_metadata(path)
        result.notes.extend(extract_notes)
        result.metadata_identifiers = metadata_terms

        terms, inferred = build_terms_for_file(
            source=path,
            extracted_text=text,
            roster_terms=roster_terms,
            metadata_terms=metadata_terms,
            auto_infer=auto_infer,
        )
        result.inferred_identifiers = inferred
        result.identifiers_found = find_terms_in_text(text + "\n" + "\n".join(metadata_terms), terms)

        if dry_run:
            result.status = "dry-run"
            return result

        if ext == ".pdf":
            removed, notes = anonymize_pdf(path, output_file, anon_id, terms)
        elif ext == ".docx":
            removed, notes = anonymize_docx(path, output_file, anon_id, terms)
        elif ext == ".pptx":
            removed, notes = anonymize_pptx(path, output_file, anon_id, terms)
        else:
            removed, notes = [], ["Unsupported file type."]

        result.identifiers_removed = removed
        result.notes.extend(notes)
        result.status = "ok"
        return result

    except Exception as e:
        result.status = "error"
        result.notes.append(str(e))
        return result


def write_manifest(results: list[FileResult], manifest_path: Path) -> None:
    manifest_path.parent.mkdir(parents=True, exist_ok=True)
    fields = [
        "anon_id",
        "source_file",
        "output_file",
        "status",
        "source_sha256",
        "identifiers_found",
        "identifiers_removed",
        "inferred_identifiers",
        "metadata_identifiers",
        "notes",
    ]

    with manifest_path.open("w", encoding="utf-8", newline="") as f:
        writer = csv.DictWriter(f, fieldnames=fields)
        writer.writeheader()
        for r in results:
            try:
                digest = sha256_file(r.source_file) if r.source_file.exists() else ""
            except Exception:
                digest = ""

            writer.writerow({
                "anon_id": r.anon_id,
                "source_file": str(r.source_file),
                "output_file": str(r.output_file or ""),
                "status": r.status,
                "source_sha256": digest,
                "identifiers_found": "; ".join(casefold_unique(r.identifiers_found)),
                "identifiers_removed": "; ".join(casefold_unique(r.identifiers_removed)),
                "inferred_identifiers": "; ".join(casefold_unique(r.inferred_identifiers)),
                "metadata_identifiers": "; ".join(casefold_unique(r.metadata_identifiers)),
                "notes": " | ".join(r.notes),
            })


def collect_files(input_dir: Path, recursive: bool) -> list[Path]:
    globber = input_dir.rglob if recursive else input_dir.glob
    files = [p for p in globber("*") if p.is_file() and not p.name.startswith(".")]
    files = [p for p in files if p.suffix.lower() in SUPPORTED_EXTS | LEGACY_OFFICE_EXTS]
    return sorted(files, key=natural_sort_key)


def main(argv: Optional[list[str]] = None) -> int:
    parser = argparse.ArgumentParser(
        description="Anonymize student PDFs, Word docs, and PowerPoints with anonymous IDs."
    )
    parser.add_argument("input_dir", type=Path, help="Directory containing student submissions.")
    parser.add_argument("-o", "--output", type=Path, default=Path("anonymized_output"),
                        help="Directory for anonymized copies and manifest.")
    parser.add_argument("--names", type=Path, default=None,
                        help="Optional roster/name list: CSV or plain text.")
    parser.add_argument("--manifest", type=Path, default=None,
                        help="Path for private crosswalk CSV. Default: output/anonymization_crosswalk.csv")
    parser.add_argument("--prefix", default="S", help="Anonymous ID prefix. Default: S")
    parser.add_argument("--start", type=int, default=1, help="Starting ID number. Default: 1")
    parser.add_argument("--width", type=int, default=4, help="Numeric width. Default: 4, e.g., S0001")
    parser.add_argument("--recursive", action="store_true", help="Process files recursively.")
    parser.add_argument("--no-auto-infer", action="store_true",
                        help="Do not infer names from filenames/text; only use --names and metadata.")
    parser.add_argument("--dry-run", action="store_true",
                        help="Create only the manifest; do not write anonymized copies.")
    args = parser.parse_args(argv)

    input_dir = args.input_dir.expanduser().resolve()
    output_dir = args.output.expanduser().resolve()
    manifest_path = (args.manifest.expanduser().resolve()
                     if args.manifest else output_dir / "anonymization_crosswalk.csv")

    if not input_dir.exists() or not input_dir.is_dir():
        print(f"Input directory does not exist or is not a directory: {input_dir}", file=sys.stderr)
        return 2

    roster_records = load_roster(args.names.expanduser().resolve() if args.names else None)
    roster_terms = all_roster_terms(roster_records)

    files = collect_files(input_dir, recursive=args.recursive)
    if not files:
        print("No supported files found.", file=sys.stderr)
        return 1

    output_dir.mkdir(parents=True, exist_ok=True)

    results: list[FileResult] = []
    for idx, path in enumerate(files, start=args.start):
        anon_id = safe_filename_id(idx, args.prefix, args.width)
        print(f"{anon_id}: {path.name}")
        results.append(process_file(
            path=path,
            anon_id=anon_id,
            output_dir=output_dir,
            roster_terms=roster_terms,
            auto_infer=not args.no_auto_infer,
            dry_run=args.dry_run,
        ))

    write_manifest(results, manifest_path)

    ok = sum(1 for r in results if r.status == "ok")
    dry = sum(1 for r in results if r.status == "dry-run")
    skipped = sum(1 for r in results if r.status == "skipped")
    errors = sum(1 for r in results if r.status == "error")

    print()
    print(f"Manifest written to: {manifest_path}")
    print(f"Processed: {ok}; dry-run: {dry}; skipped: {skipped}; errors: {errors}")
    print("Keep the manifest private; it is the re-identification crosswalk.")

    if any(r.notes for r in results):
        print()
        print("Notes/warnings were recorded in the manifest. Review before sharing anonymized files.")

    return 0 if errors == 0 else 1


if __name__ == "__main__":
    raise SystemExit(main())
